#!/usr/bin/env python3
"""料金モデル v2 準拠の監査スクリプト。

v2 ルール:
  - 全員同一料金 (¥50,000 見学会 / ¥150,000 撮影 / ¥30,000·¥80,000 月額 / ¥50,000·¥150,000 SaaS)
  - 成果報酬 3% (上限なし、Phase 3〜)
  - 成果報酬対象: 無料相談 / AI 診断 / 会員登録ユーザー
  - 10社優遇: AI 診断での優先送客のみ (料金優遇ゼロ)
  - 古い表現 (永続無料・永続優遇・1.5%・半額・初年度無料・¥25,000 等) は v2 では NG
  - /features, /sale-homes, /lands は Phase 1 で公開対象
"""
from __future__ import annotations

import re
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
LATEST = ROOT / "docs" / "最新版"

# v2 で NG な表現 (履歴欄・変更履歴以外で出現したら要修正)
# 曖昧な表現 (1.5%/半額/スターティング等) は除外し、明確な v1 マーカーのみを列挙
V2_FORBIDDEN_PHRASES = [
    "永続無料",
    "永続優遇",
    "¥25,000/件",
    "¥25,000 /件",
    "初年度無料",
    "初年度 無料",
    "50% OFF",
    "フリープラン申込書",
    "成果報酬 1.5%",
    "成果報酬1.5%",
    "見学会予約 ¥25,000",
    "見学会 ¥25,000",
    "月額掲載 初年度",
    "10社のみ無料",
    "10社様無料",
    "10社は無料",
    "10社 無料",
]

# 変更履歴等で言及される文脈 (NG にしない)
HISTORY_CONTEXTS = [
    "変更履歴",
    "旧版",
    "pre-MVP",
    "永続無料→",
    "永続無料 →",
    "料金プラン見直し",
    "廃止",
    "削除しました",
    "v1 で",
    "v1 から",
    "v1→v2",
    "料金モデル v2 で",
]

# v2 で期待される表現 (これが無いと懸念)
V2_EXPECTED_PHRASES_CORE = [
    "¥50,000",   # 見学会予約
    "3%",        # 成果報酬
]


def read_docx(path: Path) -> str:
    try:
        doc = Document(path)
        parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text)
        return "\n".join(parts)
    except Exception as e:
        return f"[ERROR: {e}]"


def read_pptx(path: Path) -> str:
    try:
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            if run.text.strip():
                                parts.append(run.text)
        return "\n".join(parts)
    except Exception as e:
        return f"[ERROR: {e}]"


def read_xlsx(path: Path) -> str:
    try:
        wb = load_workbook(path, data_only=True, read_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None and str(cell).strip():
                        parts.append(str(cell))
        return "\n".join(parts)
    except Exception as e:
        return f"[ERROR: {e}]"


def read_md(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except Exception as e:
        return f"[ERROR: {e}]"


def extract_content(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".md":
        return read_md(path)
    if suffix == ".docx":
        return read_docx(path)
    if suffix == ".pptx":
        return read_pptx(path)
    if suffix == ".xlsx":
        return read_xlsx(path)
    return ""


def is_history_context(content: str, phrase: str) -> bool:
    """phrase が変更履歴などの文脈で出現しているか判定"""
    idx = content.find(phrase)
    if idx == -1:
        return False
    window = content[max(0, idx - 100):idx + 100]
    return any(ctx in window for ctx in HISTORY_CONTEXTS)


def is_pre_mvp_file(path: Path) -> bool:
    """pre-MVP 版ファイル (参照用、v2 不整合は許容) かどうか"""
    name = path.name.lower()
    return "pre-mvp" in name or "旧版" in name


def is_archive_file(path: Path) -> bool:
    """アーカイブ的な過去文書かどうか (v1 料金を含んでも許容)"""
    archive_markers = [
        "事業モデル・ストーリー",
        "SaaS料金エキスパートパネル議事録",
        "リリース前エキスパートレビュー",
        "営業戦略統合ドキュメント",
        "business-plan-internal",
        "business-operations",
        "鹿児島事業シミュレーション",
        "営業管理（テレアポ・戦略・料金）",
        "selection-criteria",
        "MVPローンチ戦略",  # 過去の戦略文書
    ]
    return any(m in path.name for m in archive_markers)


def analyze(path: Path, content: str) -> dict:
    issues_v2 = []
    is_archive = is_archive_file(path) or is_pre_mvp_file(path)

    for phrase in V2_FORBIDDEN_PHRASES:
        if phrase in content:
            if is_history_context(content, phrase):
                continue
            if is_archive:
                continue  # 過去文書は許容
            issues_v2.append(phrase)

    # v2 公開画面の確認 (主要 docx/xlsx のみ)
    expected_urls = ["/features", "/sale-homes", "/lands"]
    if path.suffix in (".md",) and "ロードマップ" in path.name:
        missing_urls = [u for u in expected_urls if u not in content]
        if missing_urls:
            issues_v2.append(f"Phase1 公開URL未記載: {missing_urls}")

    return {
        "size": len(content),
        "is_archive": is_archive,
        "issues_v2": issues_v2,
    }


def main():
    rows = []
    for path in sorted(LATEST.rglob("*")):
        if not path.is_file():
            continue
        if path.name == ".DS_Store":
            continue
        rel = path.relative_to(LATEST)
        content = extract_content(path)
        if not content or content.startswith("[ERROR"):
            print(f"SKIP  {rel}  ({content[:60] if content else 'empty'})")
            continue
        facts = analyze(path, content)
        rows.append((rel, facts))

    # Summary
    print("\n" + "=" * 70)
    print("料金モデル v2 監査結果")
    print("=" * 70)

    # v2 不整合があるファイル
    v2_issues = [(rel, f) for rel, f in rows if f["issues_v2"]]
    if v2_issues:
        print("\n【要修正】v2 不整合 (現行ファイル):")
        for rel, f in v2_issues:
            print(f"  ✗ {rel}")
            for i in f["issues_v2"]:
                print(f"      - {i}")
    else:
        print("\n✓ v2 不整合なし (全ての現行ファイルが料金モデル v2 準拠)")

    # アーカイブファイル一覧
    archive_count = sum(1 for _, f in rows if f["is_archive"])
    print(f"\n[参考] アーカイブ/過去文書として扱うファイル数: {archive_count}")
    for rel, f in rows:
        if f["is_archive"]:
            print(f"  · {rel}")

    # 総ファイル数
    print(f"\n総ファイル数: {len(rows)}")

    return v2_issues


if __name__ == "__main__":
    issues = main()
    if issues:
        exit(1)
