#!/usr/bin/env python3
"""最新版フォルダの全ファイルから重要な事実を抽出し、重複・不整合を検出する。"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
LATEST = ROOT / "docs" / "最新版"

# 検出したいキーフレーズ
KEY_PHRASES = {
    "¥50,000": "見学会予約正規",
    "¥25,000": "見学会予約優遇",
    "¥150,000": "撮影正規",
    "¥75,000": "撮影優遇",
    "¥30,000": "月額掲載正規",
    "¥15,000": "月額掲載優遇",
    "¥50,000/月": "SaaS正規",
    "永続無料": "【要修正】永続無料",
    "50% OFF": "永続優遇",
    "2026/05/01": "リリース日",
    "2026-05-01": "リリース日",
    "Phase 1": "Phase1",
    "Phase 2": "Phase2",
    "Phase 3": "Phase3",
    "Phase 4": "Phase4",
    "Phase 5": "Phase5",
    "10社": "既存10社",
    "MVP": "MVP",
}


def read_docx(path: Path) -> str:
    try:
        doc = Document(path)
        parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text)
        return "\n".join(parts)
    except Exception as e:
        return f"[ERROR: {e}]"


def read_pptx(path: Path) -> str:
    try:
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            if run.text.strip():
                                parts.append(run.text)
        return "\n".join(parts)
    except Exception as e:
        return f"[ERROR: {e}]"


def read_xlsx(path: Path) -> str:
    try:
        wb = load_workbook(path, data_only=True, read_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None and str(cell).strip():
                        parts.append(str(cell))
        return "\n".join(parts)
    except Exception as e:
        return f"[ERROR: {e}]"


def read_md(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except Exception as e:
        return f"[ERROR: {e}]"


def extract_content(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".md":
        return read_md(path)
    if suffix == ".docx":
        return read_docx(path)
    if suffix == ".pptx":
        return read_pptx(path)
    if suffix == ".xlsx":
        return read_xlsx(path)
    return ""


def analyze(path: Path, content: str) -> dict:
    facts = {}
    for phrase, label in KEY_PHRASES.items():
        count = content.count(phrase)
        if count > 0:
            facts[label] = count

    # 日付検出 (YYYY-MM-DD または YYYY/MM/DD)
    dates = set(re.findall(r"\d{4}[-/]\d{2}[-/]\d{2}", content))
    facts["dates"] = sorted(dates)

    # 旧料金・問題のある文言を検出
    issues = []
    if "永続無料" in content:
        issues.append("永続無料")
    if "初期は無料" in content and "MVP" not in content:
        issues.append("初期は無料 (文脈要確認)")
    if "成果報酬 3%" in content and "1.5%" not in content and "MVP" not in path.name:
        issues.append("成果報酬3%のみ (1.5%併記なし)")
    if re.search(r"2025[-/]", content):
        issues.append("2025年表記あり")

    facts["issues"] = issues
    facts["size"] = len(content)
    return facts


def main():
    rows = []
    for path in sorted(LATEST.rglob("*")):
        if not path.is_file():
            continue
        if path.name == ".DS_Store":
            continue
        rel = path.relative_to(LATEST)
        content = extract_content(path)
        if not content or content.startswith("[ERROR"):
            print(f"SKIP  {rel}  ({content})")
            continue
        facts = analyze(path, content)
        rows.append((rel, facts))

    # Print summary
    print("\n═══ ファイル別サマリー ═══\n")
    for rel, facts in rows:
        print(f"■ {rel}  ({facts['size']:,} chars)")
        for key, val in facts.items():
            if key in ("dates", "issues", "size"):
                continue
            print(f"   {key}: {val}")
        if facts["dates"]:
            print(f"   dates: {', '.join(facts['dates'])}")
        if facts["issues"]:
            print(f"   ⚠️  ISSUES: {facts['issues']}")
        print()

    # Cross-file analysis
    print("\n═══ 不整合チェック ═══\n")

    # 全体で永続無料を含むファイルを列挙
    print("◆ 『永続無料』を含むファイル:")
    for rel, facts in rows:
        if "【要修正】永続無料" in facts:
            print(f"  - {rel}")

    # 2025 表記
    print("\n◆ 2025 年表記を含むファイル:")
    for rel, facts in rows:
        for issue in facts["issues"]:
            if "2025" in issue:
                print(f"  - {rel}")

    # 日付のばらつき
    print("\n◆ 検出された日付 (上位30):")
    all_dates = set()
    for rel, facts in rows:
        for d in facts["dates"]:
            all_dates.add(d)
    for d in sorted(all_dates)[:30]:
        files = [str(r) for r, f in rows if d in f["dates"]]
        print(f"  {d}: {len(files)} files")


if __name__ == "__main__":
    main()
