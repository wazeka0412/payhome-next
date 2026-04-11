#!/usr/bin/env python3
"""最新版の docx/pptx の古い日付 (2026-04-11) を 2026-04-12 に統一する。"""
from __future__ import annotations

from pathlib import Path

from docx import Document
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
LATEST = ROOT / "docs" / "最新版"

OLD = "2026-04-11"
NEW = "2026-04-12"


def fix_docx(path: Path) -> int:
    doc = Document(path)
    changed = 0
    for para in doc.paragraphs:
        for run in para.runs:
            if OLD in run.text:
                run.text = run.text.replace(OLD, NEW)
                changed += 1
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        if OLD in run.text:
                            run.text = run.text.replace(OLD, NEW)
                            changed += 1
    if changed:
        doc.save(path)
    return changed


def fix_pptx(path: Path) -> int:
    prs = Presentation(path)
    changed = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if OLD in run.text:
                        run.text = run.text.replace(OLD, NEW)
                        changed += 1
    if changed:
        prs.save(path)
    return changed


def main():
    # Only update MVP-era files that should show the latest date.
    # Skip older ops rules (03/22) and archives (which retain their original dates).
    targets = [
        LATEST / "README.docx",
        LATEST / "02_管理シート" / "00_管理シートの使い方.docx",
        LATEST / "03_営業関係" / "00_README.docx",
        LATEST / "03_営業関係" / "工務店FAQ_MVP版.docx",
        LATEST / "03_営業関係" / "料金プラン_MVP版.docx",
        LATEST / "03_営業関係" / "工務店向け説明_MVP版.pptx",
        LATEST / "04_経理関係" / "00_README.docx",
        LATEST / "04_経理関係" / "キャッシュポイント一覧.docx",
        LATEST / "05_契約書" / "00_README.docx",
        LATEST / "06_事業計画・戦略" / "00_README.docx",
        LATEST / "06_事業計画・戦略" / "MVPローンチ戦略.docx",
        LATEST / "07_業務運用ルール" / "00_README.docx",
    ]
    for path in targets:
        if not path.exists():
            print(f"skip (not found): {path.relative_to(LATEST)}")
            continue
        if path.suffix == ".docx":
            n = fix_docx(path)
        elif path.suffix == ".pptx":
            n = fix_pptx(path)
        else:
            continue
        print(f"{'fixed' if n else 'no-change':>10} {n:>3}x  {path.relative_to(LATEST)}")


if __name__ == "__main__":
    main()
