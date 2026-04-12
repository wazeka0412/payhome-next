#!/usr/bin/env python3
"""最新版フォルダの .md ファイルを .docx / .pptx に変換する。

方針:
  - 01_エンジニア共有用/ は .md のまま保持
  - 02_管理シート/, 03_営業関係/, 04_経理関係/, 最新版/README は .docx / .pptx 化
  - 工務店向け説明_MVP版.md のみ 説明スライドとして .pptx 化
  - それ以外は .docx 化
"""
from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PPTXColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PPTXInches, Pt as PPTXPt

ROOT = Path(__file__).resolve().parent.parent
LATEST = ROOT / "docs" / "最新版"

# ─── Brand colors ─────────────────────────────────────────────────

BROWN = RGBColor(0x3D, 0x22, 0x00)
ORANGE = RGBColor(0xE8, 0x74, 0x0C)
GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK = RGBColor(0x11, 0x18, 0x27)

PPTX_BROWN = PPTXColor(0x3D, 0x22, 0x00)
PPTX_ORANGE = PPTXColor(0xE8, 0x74, 0x0C)
PPTX_CREAM = PPTXColor(0xFF, 0xF8, 0xF0)
PPTX_WHITE = PPTXColor(0xFF, 0xFF, 0xFF)
PPTX_GRAY = PPTXColor(0x6B, 0x72, 0x80)


# ══════════════════════════════════════════════════════════════════
# Markdown parser (minimal)
# ══════════════════════════════════════════════════════════════════

def parse_md(text: str) -> list[dict]:
    """Parse markdown into a list of blocks. Each block is {type, ...}."""
    blocks: list[dict] = []
    lines = text.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # Heading
        if stripped.startswith("#"):
            level = len(stripped) - len(stripped.lstrip("#"))
            content = stripped.lstrip("#").strip()
            blocks.append({"type": "heading", "level": level, "text": content})
            i += 1
            continue

        # Horizontal rule
        if stripped in {"---", "***"}:
            blocks.append({"type": "hr"})
            i += 1
            continue

        # Blank line
        if not stripped:
            blocks.append({"type": "blank"})
            i += 1
            continue

        # Table
        if stripped.startswith("|"):
            table_rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                row = [cell.strip() for cell in lines[i].strip().strip("|").split("|")]
                table_rows.append(row)
                i += 1
            # Drop separator row like |---|---|
            table_rows = [r for r in table_rows if not all(set(c) <= set("- :") for c in r)]
            if table_rows:
                blocks.append({"type": "table", "rows": table_rows})
            continue

        # List
        if stripped.startswith(("- ", "* ", "+ ")):
            items = []
            while i < len(lines) and lines[i].strip().startswith(("- ", "* ", "+ ")):
                items.append(lines[i].strip()[2:])
                i += 1
            blocks.append({"type": "list", "ordered": False, "items": items})
            continue

        # Numbered list
        if stripped and stripped[0].isdigit() and ". " in stripped[:5]:
            items = []
            while i < len(lines):
                s = lines[i].strip()
                if s and s[0].isdigit() and ". " in s[:5]:
                    items.append(s[s.index(". ") + 2 :])
                    i += 1
                else:
                    break
            blocks.append({"type": "list", "ordered": True, "items": items})
            continue

        # Code fence
        if stripped.startswith("```"):
            i += 1
            code_lines = []
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1  # skip closing fence
            blocks.append({"type": "code", "text": "\n".join(code_lines)})
            continue

        # Paragraph
        para_lines = [line]
        i += 1
        while i < len(lines):
            nxt = lines[i]
            if not nxt.strip() or nxt.strip().startswith(("#", "|", "- ", "```", "---")):
                break
            para_lines.append(nxt)
            i += 1
        blocks.append({"type": "para", "text": " ".join(p.strip() for p in para_lines)})

    return blocks


# ══════════════════════════════════════════════════════════════════
# DOCX writer
# ══════════════════════════════════════════════════════════════════

def blocks_to_docx(blocks: list[dict], title: str, out_path: Path) -> None:
    doc = Document()

    # Set default font
    style = doc.styles["Normal"]
    style.font.name = "Hiragino Sans"
    style.font.size = Pt(10)

    # Document title
    h = doc.add_heading(title, level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for run in h.runs:
        run.font.color.rgb = BROWN
        run.font.size = Pt(22)

    for block in blocks:
        btype = block["type"]

        if btype == "heading":
            level = min(block["level"], 4)
            h = doc.add_heading(block["text"], level=level)
            for run in h.runs:
                run.font.color.rgb = BROWN if level <= 2 else ORANGE
                run.font.size = Pt(18 - level * 2)

        elif btype == "para":
            p = doc.add_paragraph()
            text = block["text"]
            # Simple bold handling
            parts = _split_bold(text)
            for t, bold in parts:
                run = p.add_run(t)
                run.bold = bold
                run.font.size = Pt(10)
                run.font.color.rgb = DARK

        elif btype == "list":
            for item in block["items"]:
                p = doc.add_paragraph(style="List Bullet" if not block["ordered"] else "List Number")
                parts = _split_bold(item)
                for t, bold in parts:
                    run = p.add_run(t)
                    run.bold = bold
                    run.font.size = Pt(10)

        elif btype == "table":
            rows = block["rows"]
            if not rows:
                continue
            table = doc.add_table(rows=len(rows), cols=len(rows[0]))
            table.style = "Light Grid Accent 1"
            for r_idx, row in enumerate(rows):
                for c_idx, cell_text in enumerate(row):
                    if c_idx < len(table.rows[r_idx].cells):
                        cell = table.rows[r_idx].cells[c_idx]
                        cell.text = ""
                        p = cell.paragraphs[0]
                        parts = _split_bold(cell_text)
                        for t, bold in parts:
                            run = p.add_run(t)
                            run.bold = bold or r_idx == 0
                            run.font.size = Pt(9)
                            if r_idx == 0:
                                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        if r_idx == 0:
                            _shade_cell(cell, "3D2200")

        elif btype == "code":
            p = doc.add_paragraph()
            run = p.add_run(block["text"])
            run.font.name = "Menlo"
            run.font.size = Pt(8)
            run.font.color.rgb = DARK
            _shade_paragraph(p, "F3F4F6")

        elif btype == "hr":
            doc.add_paragraph("─" * 50).alignment = WD_ALIGN_PARAGRAPH.CENTER

        elif btype == "blank":
            doc.add_paragraph("")

    doc.save(out_path)


def _split_bold(text: str) -> list[tuple[str, bool]]:
    """Very small ``**bold**`` splitter."""
    parts: list[tuple[str, bool]] = []
    buf = ""
    i = 0
    while i < len(text):
        if text[i : i + 2] == "**":
            if buf:
                parts.append((buf, False))
                buf = ""
            j = text.find("**", i + 2)
            if j == -1:
                buf += text[i:]
                break
            parts.append((text[i + 2 : j], True))
            i = j + 2
        else:
            buf += text[i]
            i += 1
    if buf:
        parts.append((buf, False))
    return parts


def _shade_cell(cell, color_hex: str) -> None:
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), color_hex)
    tc_pr.append(shd)


def _shade_paragraph(p, color_hex: str) -> None:
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    pPr = p._p.get_or_add_pPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), color_hex)
    pPr.append(shd)


# ══════════════════════════════════════════════════════════════════
# PPTX writer (for the 工務店 explanation slide deck)
# ══════════════════════════════════════════════════════════════════

def blocks_to_pptx(blocks: list[dict], title: str, subtitle: str, out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = PPTXInches(13.333)
    prs.slide_height = PPTXInches(7.5)

    # ── Title slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _fill_background(slide, PPTX_BROWN)
    _add_text(
        slide, 0.5, 2.5, 12.3, 1.5,
        title, size=36, bold=True, color=PPTX_WHITE, align=PP_ALIGN.LEFT,
    )
    _add_text(
        slide, 0.5, 4.0, 12.3, 0.6,
        subtitle, size=18, bold=False, color=PPTX_CREAM, align=PP_ALIGN.LEFT,
    )
    _add_text(
        slide, 0.5, 6.5, 12.3, 0.4,
        "株式会社wazeka / ぺいほーむ事業部",
        size=12, bold=False, color=PPTX_CREAM, align=PP_ALIGN.LEFT,
    )

    # ── Convert H2 sections into section-header + content slides ──
    current_h2: str | None = None
    current_content: list[dict] = []

    def flush() -> None:
        if not current_h2:
            return
        # Split content into chunks that fit on a slide (~8 items each)
        chunks = _chunk_content(current_content, max_items=8)
        for idx, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_rect(slide, 0, 0, 13.333, 0.8, PPTX_ORANGE)
            heading = f"{current_h2}" + (f"  ({idx+1}/{len(chunks)})" if len(chunks) > 1 else "")
            _add_text(
                slide, 0.3, 0.1, 13.0, 0.6,
                heading, size=22, bold=True, color=PPTX_WHITE, align=PP_ALIGN.LEFT,
            )
            _render_chunk(slide, chunk, y_start=1.2)

    for block in blocks:
        if block["type"] == "heading" and block["level"] == 2:
            flush()
            current_h2 = block["text"]
            current_content = []
        elif block["type"] == "heading" and block["level"] == 1:
            # Top-level heading -> treat as subtitle on title slide, skip
            continue
        elif current_h2 is not None:
            current_content.append(block)

    flush()

    # Closing slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_background(slide, PPTX_BROWN)
    _add_text(
        slide, 0.5, 2.8, 12.3, 1.2,
        "よろしくお願いいたします。",
        size=40, bold=True, color=PPTX_WHITE, align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide, 0.5, 4.2, 12.3, 0.6,
        "support@payhome.jp",
        size=18, bold=False, color=PPTX_CREAM, align=PP_ALIGN.CENTER,
    )

    prs.save(out_path)


def _chunk_content(blocks: list[dict], max_items: int = 8) -> list[list[dict]]:
    chunks = []
    current: list[dict] = []
    count = 0
    for b in blocks:
        if b["type"] == "blank":
            continue
        # Count items in this block
        if b["type"] == "list":
            n = len(b["items"])
        elif b["type"] == "table":
            n = len(b["rows"])
        elif b["type"] == "heading":
            n = 1
        elif b["type"] == "para":
            n = 2
        else:
            n = 1

        if count + n > max_items and current:
            chunks.append(current)
            current = [b]
            count = n
        else:
            current.append(b)
            count += n
    if current:
        chunks.append(current)
    return chunks or [[]]


def _render_chunk(slide, blocks: list[dict], y_start: float) -> None:
    y = y_start
    for block in blocks:
        btype = block["type"]

        if btype == "heading":
            level = block["level"]
            _add_text(
                slide, 0.4, y, 12.5, 0.5,
                block["text"],
                size=18 if level == 3 else 14,
                bold=True,
                color=PPTX_BROWN,
            )
            y += 0.55

        elif btype == "para":
            _add_text(
                slide, 0.5, y, 12.4, 0.5,
                block["text"],
                size=13, bold=False, color=PPTX_BROWN,
            )
            y += 0.55

        elif btype == "list":
            for item in block["items"]:
                _add_text(
                    slide, 0.7, y, 12.2, 0.42,
                    "• " + item,
                    size=12, bold=False, color=PPTX_BROWN,
                )
                y += 0.42

        elif btype == "table":
            rows = block["rows"]
            if not rows:
                continue
            cols = len(rows[0])
            col_w = 12.4 / cols
            row_h = 0.42
            for r_idx, row in enumerate(rows):
                for c_idx, cell in enumerate(row):
                    x = 0.4 + c_idx * col_w
                    _add_text(
                        slide, x, y, col_w, row_h,
                        cell,
                        size=10, bold=(r_idx == 0), color=PPTX_BROWN,
                    )
                    if r_idx == 0:
                        _add_rect(slide, x, y, col_w, row_h, PPTX_CREAM, behind=True)
                y += row_h

        elif btype == "hr":
            y += 0.3

        if y > 7.0:
            break


def _add_text(slide, x: float, y: float, w: float, h: float, text: str,
              *, size: int = 12, bold: bool = False, color=PPTX_BROWN,
              align=PP_ALIGN.LEFT) -> None:
    tb = slide.shapes.add_textbox(PPTXInches(x), PPTXInches(y),
                                  PPTXInches(w), PPTXInches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = PPTXPt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Hiragino Sans"


def _add_rect(slide, x: float, y: float, w: float, h: float, color, behind: bool = False) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        PPTXInches(x), PPTXInches(y),
        PPTXInches(w), PPTXInches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if behind:
        # Move behind all other shapes
        spTree = shape._element.getparent()
        spTree.remove(shape._element)
        spTree.insert(2, shape._element)


def _fill_background(slide, color) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        PPTXInches(13.333), PPTXInches(7.5),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # Move to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


# ══════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════

def convert_md_to_docx(md_path: Path, docx_path: Path, title: str | None = None) -> None:
    text = md_path.read_text(encoding="utf-8")
    blocks = parse_md(text)

    if title is None:
        # Use first H1 or filename
        first_h1 = next(
            (b["text"] for b in blocks if b["type"] == "heading" and b["level"] == 1),
            None,
        )
        title = first_h1 or md_path.stem

    # Remove the first H1 from blocks to avoid duplication
    blocks = [b for b in blocks if not (b["type"] == "heading" and b["level"] == 1)]

    blocks_to_docx(blocks, title, docx_path)


def convert_md_to_pptx(md_path: Path, pptx_path: Path, subtitle: str) -> None:
    text = md_path.read_text(encoding="utf-8")
    blocks = parse_md(text)

    first_h1 = next(
        (b["text"] for b in blocks if b["type"] == "heading" and b["level"] == 1),
        md_path.stem,
    )
    blocks = [b for b in blocks if not (b["type"] == "heading" and b["level"] == 1)]

    blocks_to_pptx(blocks, first_h1, subtitle, pptx_path)


def main() -> None:
    tasks = [
        # (source md, target, subtitle for pptx, format)
        (LATEST / "README.md", LATEST / "README.docx", None, "docx"),
        (LATEST / "02_管理シート" / "00_管理シートの使い方.md",
         LATEST / "02_管理シート" / "00_管理シートの使い方.docx",
         None, "docx"),
        (LATEST / "03_営業関係" / "00_README.md",
         LATEST / "03_営業関係" / "00_README.docx",
         None, "docx"),
        (LATEST / "03_営業関係" / "工務店向け説明_MVP版.md",
         LATEST / "03_営業関係" / "工務店向け説明_MVP版.pptx",
         "2026-05-01 MVPリリース", "pptx"),
        (LATEST / "03_営業関係" / "料金プラン_MVP版.md",
         LATEST / "03_営業関係" / "料金プラン_MVP版.docx",
         None, "docx"),
        (LATEST / "03_営業関係" / "工務店FAQ_MVP版.md",
         LATEST / "03_営業関係" / "工務店FAQ_MVP版.docx",
         None, "docx"),
        (LATEST / "04_経理関係" / "00_README.md",
         LATEST / "04_経理関係" / "00_README.docx",
         None, "docx"),
        (LATEST / "04_経理関係" / "キャッシュポイント一覧.md",
         LATEST / "04_経理関係" / "キャッシュポイント一覧.docx",
         None, "docx"),
    ]

    for src, dst, subtitle, fmt in tasks:
        if not src.exists():
            print(f"skip (not found): {src}")
            continue
        if fmt == "docx":
            convert_md_to_docx(src, dst)
        else:
            convert_md_to_pptx(src, dst, subtitle or "")
        print(f"created: {dst}")
        # Remove source
        src.unlink()
        print(f"removed: {src}")


if __name__ == "__main__":
    main()
